"""
Script to create comprehensive PowerPoint presentation for Final Report
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(37, 99, 235)  # Blue
    text_color = RGBColor(31, 41, 55)    # Dark gray
    accent_color = RGBColor(16, 185, 129)  # Green
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = text_color
        return slide
    
    def add_content_slide(title, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, point in enumerate(bullet_points):
            p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = text_color
            p.level = 0
            p.space_after = Pt(12)
        
        return slide
    
    def add_two_column_slide(title, left_points, right_points):
        slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two content layout
        title_shape = slide.shapes.title
        
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Left column
        left_shape = slide.placeholders[1]
        left_frame = left_shape.text_frame
        left_frame.clear()
        for point in left_points:
            p = left_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.level = 0
        
        # Right column
        right_shape = slide.placeholders[2]
        right_frame = right_shape.text_frame
        right_frame.clear()
        for point in right_points:
            p = right_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.level = 0
        
        return slide
    
    # Slide 1: Title Slide
    add_title_slide(
        "Machine Learning for Burnout Risk Prediction",
        "A Comprehensive Solution for Hybrid and Remote Teams\nITC-586 Design Studio | Fall 2025"
    )
    
    # Slide 2: Agenda
    add_content_slide("Agenda", [
        "Problem Statement & Objectives",
        "System Architecture & Design",
        "Machine Learning Approach",
        "Implementation Highlights",
        "Results & Performance",
        "Risk Management",
        "Conclusions & Recommendations",
        "Reflection & Lessons Learned"
    ])
    
    # Slide 3: Problem Statement
    add_content_slide("Problem Statement", [
        "Burnout affects 23% of employees globally",
        "Hybrid/remote work creates new challenges for monitoring well-being",
        "Traditional detection methods are reactive, not proactive",
        "Organizations lack real-time, data-driven burnout prediction tools",
        "Need for early intervention to prevent critical burnout"
    ])
    
    # Slide 4: Objectives
    add_two_column_slide(
        "Project Objectives",
        [
            "Develop accurate ML prediction model",
            "Integrate multiple data sources",
            "Create user-friendly interface",
            "Implement enterprise security"
        ],
        [
            "Achieve 80%+ test coverage",
            "Support 1000+ concurrent users",
            "Enable real-time predictions",
            "Provide actionable insights"
        ]
    )
    
    # Slide 5: System Architecture
    add_content_slide("System Architecture", [
        "Microservices Architecture: Frontend, Backend, ML Service",
        "Frontend: React + TypeScript (Port 5173)",
        "Backend: Node.js + Express (Port 3001)",
        "ML Service: Python + FastAPI (Port 8000)",
        "Database: MongoDB (Port 27017)",
        "Cache: Redis (Port 6379)"
    ])
    
    # Slide 6: Architecture Diagram (Text representation)
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture Overview"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = title_color
    
    arch_text = """┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│  Frontend   │    │   Backend   │    │ ML Service  │
│   React     │◄──►│  Node.js    │◄──►│   Python    │
│  Port 5173  │    │  Port 3001  │    │  Port 8000  │
└─────────────┘    └─────────────┘    └─────────────┘
       │                  │                  │
       │                  │                  │
       ▼                  ▼                  ▼
┌─────────────┐    ┌─────────────┐    ┌─────────────┐
│             │    │   MongoDB   │    │    Redis     │
│   (Proxy)   │    │   Database  │    │    Cache     │
└─────────────┘    └─────────────┘    └─────────────┘"""
    
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    text_frame = text_box.text_frame
    text_frame.text = arch_text
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.name = "Courier New"
    
    # Slide 7: Machine Learning Approach
    add_content_slide("Machine Learning Approach", [
        "Ensemble Method: Combines multiple models for robust predictions",
        "Baseline Models: Logistic Regression (96% accuracy)",
        "Baseline Models: Random Forest (93% accuracy)",
        "Advanced Models: BERT for text sentiment analysis",
        "Advanced Models: LSTM for sequential pattern recognition",
        "Risk Levels: Low, Medium, High, Critical"
    ])
    
    # Slide 8: Data Sources
    add_two_column_slide(
        "Data Sources",
        [
            "Calendar Events",
            "• Meeting frequency",
            "• After-hours work",
            "• Meeting duration",
            "",
            "Email Communications",
            "• Sentiment analysis",
            "• Response times",
            "• Communication patterns"
        ],
        [
            "Survey Responses",
            "• Self-reported stress",
            "• Work-life balance",
            "• Job satisfaction",
            "",
            "Behavioral Patterns",
            "• Work hours",
            "• Break patterns",
            "• Activity levels"
        ]
    )
    
    # Slide 9: Feature Engineering
    add_content_slide("Feature Engineering", [
        "Temporal Features: Work hours, meeting frequency, after-hours activity",
        "Communication Features: Email volume, response time, sentiment scores",
        "Behavioral Features: Work-life balance indicators, break patterns",
        "Survey Features: Self-reported stress, satisfaction, support levels",
        "Top Predictors: Work hours (0.68), After-hours meetings (0.62), Email sentiment (0.58)"
    ])
    
    # Slide 10: Model Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = "Model Performance Results"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    metrics = [
        "Overall Accuracy: 87.3%",
        "Macro F1-Score: 0.84",
        "ROC-AUC: 0.89",
        "",
        "Per-Model Performance:",
        "• Logistic Regression: 96% accuracy, 0.96 F1-score",
        "• Random Forest: 93% accuracy, 0.93 F1-score",
        "• Ensemble: 87.3% accuracy, 0.84 F1-score"
    ]
    
    for i, metric in enumerate(metrics):
        p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
        p.text = metric
        p.font.size = Pt(18) if not metric.startswith("•") else Pt(16)
        p.font.color.rgb = text_color
        p.level = 1 if metric.startswith("•") else 0
    
    # Slide 11: Implementation Highlights - Backend
    add_content_slide("Backend Implementation", [
        "RESTful API with Express.js and TypeScript",
        "JWT Authentication with refresh tokens",
        "Role-Based Access Control (RBAC)",
        "MongoDB integration with Mongoose ODM",
        "Redis caching for performance optimization",
        "Rate limiting and security headers",
        "Comprehensive error handling and logging"
    ])
    
    # Slide 12: Implementation Highlights - Frontend
    add_content_slide("Frontend Implementation", [
        "React 18 with TypeScript for type safety",
        "Role-based dashboards (Employee, Manager, Admin)",
        "Interactive data visualizations with Recharts",
        "Responsive design for all devices",
        "Real-time data updates",
        "Intuitive user interface",
        "Accessibility compliance (WCAG 2.1 AA)"
    ])
    
    # Slide 13: Implementation Highlights - ML Service
    add_content_slide("ML Service Implementation", [
        "FastAPI for high-performance API",
        "Automated training pipeline",
        "Model versioning and persistence",
        "EDA report generation",
        "Real-time prediction inference",
        "Feature engineering pipeline",
        "Model evaluation and metrics tracking"
    ])
    
    # Slide 14: Testing & Validation
    add_two_column_slide(
        "Testing & Validation",
        [
            "ML Tests: 80%+ coverage",
            "• Unit tests",
            "• Model validation",
            "• Cross-validation",
            "",
            "Backend Tests: 80%+ coverage",
            "• Unit tests",
            "• Integration tests",
            "• API endpoint tests"
        ],
        [
            "Frontend Tests: 80%+ coverage",
            "• Component tests",
            "• Integration tests",
            "• User flow tests",
            "",
            "System Tests",
            "• Performance testing",
            "• Security testing",
            "• End-to-end testing"
        ]
    )
    
    # Slide 15: Security Features
    add_content_slide("Security Features", [
        "JWT Authentication with secure token management",
        "Password hashing using bcrypt",
        "Role-Based Access Control (RBAC)",
        "Rate limiting on all API endpoints",
        "Input validation and sanitization",
        "CORS protection",
        "Helmet.js security headers",
        "Secrets management via environment variables"
    ])
    
    # Slide 16: Risk Management
    add_two_column_slide(
        "Risk Management",
        [
            "Technical Risks",
            "• Model accuracy",
            "• Data privacy",
            "• System scalability",
            "• Integration failures",
            "",
            "Mitigation Strategies",
            "• Ensemble methods",
            "• Encryption & access controls",
            "• Caching & optimization"
        ],
        [
            "Operational Risks",
            "• User adoption",
            "• Data quality",
            "",
            "Security Risks",
            "• Unauthorized access",
            "• Data exposure",
            "",
            "Monitoring",
            "• Health checks",
            "• Security logging",
            "• Performance metrics"
        ]
    )
    
    # Slide 17: Key Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = "Key Results & Achievements"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    results = [
        "✓ 87.3% prediction accuracy achieved",
        "✓ 80%+ test coverage across all components",
        "✓ Production-ready security implementation",
        "✓ Scalable architecture supporting 1000+ users",
        "✓ Real-time prediction capability (< 500ms)",
        "✓ Comprehensive documentation and testing",
        "✓ User-friendly interface with role-based access"
    ]
    
    for i, result in enumerate(results):
        p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
        p.text = result
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color if result.startswith("✓") else text_color
        p.level = 0
    
    # Slide 18: Conclusions
    add_content_slide("Conclusions", [
        "Successfully developed comprehensive ML system for burnout prediction",
        "Ensemble model achieved 87.3% accuracy with robust performance",
        "Production-ready system with enterprise-grade security",
        "Effective integration of multiple data sources",
        "User-centric design enables effective use by all stakeholders",
        "Scalable architecture supports future expansion",
        "System demonstrates readiness for organizational deployment"
    ])
    
    # Slide 19: Recommendations - Implementation
    add_content_slide("Recommendations: Implementation", [
        "Phased rollout starting with pilot group",
        "Comprehensive user training for all roles",
        "Continuous monitoring and performance tracking",
        "Regular model retraining (quarterly recommended)",
        "Clear privacy and ethics communication",
        "Establish feedback mechanisms for continuous improvement"
    ])
    
    # Slide 20: Recommendations - Future Enhancements
    add_content_slide("Recommendations: Future Enhancements", [
        "Integrate biometric data (sleep, activity levels)",
        "Add collaboration tool metrics (Slack, Teams)",
        "Develop automated intervention recommendations",
        "Implement longitudinal trend analysis",
        "Add explainable AI features",
        "Develop native mobile applications"
    ])
    
    # Slide 21: Lessons Learned
    add_two_column_slide(
        "Lessons Learned",
        [
            "Start simple, iterate",
            "Testing is critical",
            "Documentation matters",
            "Security first approach",
            "User feedback is valuable"
        ],
        [
            "Performance considerations",
            "Team communication",
            "Flexibility is essential",
            "Continuous learning",
            "Quality over speed"
        ]
    )
    
    # Slide 22: Challenges & Solutions
    add_content_slide("Challenges & Solutions", [
        "Challenge: Model overfitting → Solution: Cross-validation and ensemble methods",
        "Challenge: API performance → Solution: Redis caching and optimization",
        "Challenge: Data quality → Solution: Validation pipelines",
        "Challenge: Security concerns → Solution: Security reviews and best practices",
        "Challenge: Integration complexity → Solution: Clear APIs and documentation"
    ])
    
    # Slide 23: Team Reflection
    add_content_slide("Team Reflection", [
        "Strong collaboration and effective division of labor",
        "Clear role definition enabled focused contributions",
        "Regular communication kept team aligned",
        "Quality focus maintained throughout development",
        "Adaptability to challenges and changing requirements",
        "Collective commitment to delivering high-quality work"
    ])
    
    # Slide 24: Impact & Value Proposition
    add_content_slide("Impact & Value Proposition", [
        "Early detection enables proactive intervention",
        "Reduces burnout-related costs (turnover, productivity loss)",
        "Improves employee well-being and job satisfaction",
        "Provides data-driven insights for management",
        "Supports evidence-based HR decision-making",
        "Scalable solution for organizations of all sizes"
    ])
    
    # Slide 25: Next Steps
    add_content_slide("Next Steps", [
        "Pilot deployment with selected team",
        "Gather user feedback and iterate",
        "Expand data source integrations",
        "Enhance model explainability",
        "Develop intervention recommendation engine",
        "Scale to full organizational deployment"
    ])
    
    # Slide 26: Thank You
    add_title_slide(
        "Thank You",
        "Questions & Discussion\n\nMachine Learning for Burnout Risk Prediction\nITC-586 Design Studio | Fall 2025"
    )
    
    return prs

if __name__ == "__main__":
    print("Creating PowerPoint presentation...")
    prs = create_presentation()
    output_file = "Final_Presentation_Burnout_Prediction_System.pptx"
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")
    print(f"Total slides: {len(prs.slides)}")


